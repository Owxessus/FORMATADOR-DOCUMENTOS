# -*- coding: utf-8 -*-
"""Apresentações (.pptx): criar, ler e editar.

Criar é o caso comum: a IA devolve um PLANO (título, slides, tópicos,
notas) e este módulo monta o arquivo. Editar funciona por operações
nomeadas — a IA nunca executa código, igual ao que já é feito no Excel.

Cuidado tomado: ao trocar texto de uma forma existente, escrevemos no
`run` e não em `text_frame.text`, porque esse atalho apaga toda a
formatação do parágrafo.
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# paleta sóbria, adequada a documento institucional
COR_TITULO = RGBColor(0x1E, 0x35, 0x5E)
COR_TEXTO = RGBColor(0x22, 0x22, 0x22)
COR_APOIO = RGBColor(0x55, 0x66, 0x77)
FONTE = "Calibri"


# --------------------------------------------------------------- criar

def criar(plano: dict, destino: str) -> str:
    """Monta a apresentação a partir do plano.

    plano = {"titulo": "...", "subtitulo": "...",
             "slides": [{"titulo": "...", "topicos": ["..."],
                         "texto": "...", "notas": "..."}]}
    """
    pres = Presentation()
    pres.slide_width = Inches(13.333)      # 16:9
    pres.slide_height = Inches(7.5)

    if plano.get("titulo"):
        s = pres.slides.add_slide(pres.slide_layouts[0])
        _ocupar_slide(s, pres)
        s.shapes.title.text = plano["titulo"]
        _estilo(s.shapes.title.text_frame, 40, True, COR_TITULO)
        if len(s.placeholders) > 1 and plano.get("subtitulo"):
            sub = s.placeholders[1]
            sub.text = plano["subtitulo"]
            _estilo(sub.text_frame, 18, False, COR_APOIO)
        if plano.get("notas"):
            s.notes_slide.notes_text_frame.text = str(plano["notas"])

    for item in plano.get("slides", []):
        topicos = item.get("topicos") or []
        corpo = (item.get("texto") or "").strip()
        layout = pres.slide_layouts[1] if (topicos or corpo) \
            else pres.slide_layouts[5]
        s = pres.slides.add_slide(layout)
        _ocupar_slide(s, pres)
        if s.shapes.title is not None:
            s.shapes.title.text = item.get("titulo", "")
            _estilo(s.shapes.title.text_frame, 30, True, COR_TITULO)

        if topicos or corpo:
            quadro = None
            for ph in s.placeholders:
                if ph.placeholder_format.idx != 0:
                    quadro = ph.text_frame
                    break
            if quadro is None:
                quadro = s.shapes.add_textbox(
                    Inches(0.9), Inches(1.8), Inches(11.5),
                    Inches(4.8)).text_frame
            quadro.word_wrap = True
            linhas = topicos or [l for l in corpo.split("\n") if l.strip()]
            for i, linha in enumerate(linhas):
                par = quadro.paragraphs[0] if i == 0 else \
                    quadro.add_paragraph()
                par.text = str(linha)
                par.level = 0
                par.space_after = Pt(10)
                for run in par.runs:
                    run.font.size = Pt(18)
                    run.font.name = FONTE
                    run.font.color.rgb = COR_TEXTO
        if item.get("notas"):
            s.notes_slide.notes_text_frame.text = str(item["notas"])

    os.makedirs(os.path.dirname(os.path.abspath(destino)), exist_ok=True)
    pres.save(destino)
    return destino


def _ocupar_slide(s, pres):
    """Reposiciona os marcadores do modelo para a largura 16:9.

    O modelo padrão do python-pptx é 4:3; mudar slide_width não move os
    placeholders, e o conteúdo fica espremido no canto esquerdo.
    """
    larg = pres.slide_width
    margem = Inches(0.8)
    if s.shapes.title is not None:
        t = s.shapes.title
        t.left, t.top = margem, Inches(0.45)
        t.width, t.height = larg - 2 * margem, Inches(1.15)
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            continue
        ph.left, ph.top = margem, Inches(1.85)
        ph.width = larg - 2 * margem
        ph.height = pres.slide_height - Inches(2.5)


def _estilo(quadro, tamanho, negrito, cor):
    for par in quadro.paragraphs:
        for run in par.runs:
            run.font.size = Pt(tamanho)
            run.font.bold = negrito
            run.font.name = FONTE
            run.font.color.rgb = cor


# ----------------------------------------------------------------- ler

def ler(caminho: str) -> str:
    """Conteúdo da apresentação em texto, slide a slide."""
    pres = Presentation(caminho)
    saida = []
    for n, s in enumerate(pres.slides, 1):
        partes = [f"[Slide {n}]"]
        for forma in s.shapes:
            if forma.has_text_frame and forma.text_frame.text.strip():
                partes.append("  " + forma.text_frame.text.strip()
                              .replace("\n", "\n  "))
        if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip():
            partes.append("  (notas) "
                          + s.notes_slide.notes_text_frame.text.strip())
        saida.append("\n".join(partes))
    return "\n\n".join(saida)


# -------------------------------------------------------------- editar

def editar(caminho: str, operacoes: list[dict],
           destino: str | None = None) -> str:
    """Aplica operações e salva uma CÓPIA.

    Operações aceitas:
      {"tipo":"substituir_texto","de":"XXX","para":"YYY"}
      {"tipo":"alterar_titulo","slide":2,"valor":"..."}
      {"tipo":"adicionar_slide","titulo":"...","topicos":[...],"notas":"..."}
      {"tipo":"remover_slide","slide":3}
      {"tipo":"notas","slide":1,"valor":"..."}
    """
    pres = Presentation(caminho)

    for op in operacoes:
        tipo = (op.get("tipo") or "").lower()

        if tipo == "substituir_texto":
            de, para = op.get("de", ""), op.get("para", "")
            if not de:
                continue
            for s in pres.slides:
                for forma in s.shapes:
                    if not forma.has_text_frame:
                        continue
                    for par in forma.text_frame.paragraphs:
                        for run in par.runs:      # nunca text_frame.text:
                            if de in run.text:    # isso apagaria o estilo
                                run.text = run.text.replace(de, para)

        elif tipo == "alterar_titulo":
            s = _slide(pres, op.get("slide"))
            if s is not None and s.shapes.title is not None:
                tf = s.shapes.title.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = op.get("valor", "")
                    for extra in tf.paragraphs[0].runs[1:]:
                        extra.text = ""
                else:
                    s.shapes.title.text = op.get("valor", "")

        elif tipo == "adicionar_slide":
            criar_slide_em(pres, op)

        elif tipo == "remover_slide":
            s = _slide(pres, op.get("slide"))
            if s is not None:
                ident = s.slide_id
                lst = pres.slides._sldIdLst
                for el in list(lst):
                    if int(el.get(
                        "{http://schemas.openxmlformats.org/officeDocument"
                        "/2006/relationships}id", "0").replace("rId", "0")
                            or 0) or True:
                        pass
                for el, slide in zip(list(lst), list(pres.slides)):
                    if slide.slide_id == ident:
                        lst.remove(el)
                        break

        elif tipo == "notas":
            s = _slide(pres, op.get("slide"))
            if s is not None:
                s.notes_slide.notes_text_frame.text = op.get("valor", "")

    if destino is None:
        base, ext = os.path.splitext(caminho)
        destino = f"{base}_EDITADO{ext}"
    pres.save(destino)
    return destino


def _slide(pres, numero):
    try:
        n = int(numero) - 1
        return pres.slides[n] if 0 <= n < len(pres.slides) else None
    except (TypeError, ValueError):
        return None


def criar_slide_em(pres, item: dict):
    s = pres.slides.add_slide(pres.slide_layouts[1])
    _ocupar_slide(s, pres)
    if s.shapes.title is not None:
        s.shapes.title.text = item.get("titulo", "")
        _estilo(s.shapes.title.text_frame, 30, True, COR_TITULO)
    quadro = None
    for ph in s.placeholders:
        if ph.placeholder_format.idx != 0:
            quadro = ph.text_frame
            break
    if quadro is not None:
        for i, linha in enumerate(item.get("topicos") or []):
            par = quadro.paragraphs[0] if i == 0 else quadro.add_paragraph()
            par.text = str(linha)
            for run in par.runs:
                run.font.size = Pt(18)
                run.font.name = FONTE
                run.font.color.rgb = COR_TEXTO
    if item.get("notas"):
        s.notes_slide.notes_text_frame.text = str(item["notas"])
    return s
